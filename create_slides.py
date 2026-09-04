import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

def create_presentation():
    # 1. 프레젠테이션 초기화 및 16:9 슬라이드 크기 설정 (13.333 x 7.5 인치)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # 빈 슬라이드 레이아웃

    # 색상 상수 정의
    COLOR_NAVY = RGBColor(16, 44, 87)       # 제목용 남색 (Navy #102C57)
    COLOR_BLACK = RGBColor(0, 0, 0)         # 본문용 검정
    COLOR_WHITE = RGBColor(255, 255, 255)   # 코드/명령어용 흰색
    COLOR_CODE_BG = RGBColor(24, 28, 36)    # 코드 블록 배경 (Dark/Black)
    COLOR_BORDER = RGBColor(200, 210, 225)  # 구분선 및 엣지 색상

    def set_run_font(run, text, size_pt, bold=False, color_rgb=COLOR_BLACK, font_name="맑은 고딕"):
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color_rgb

        # 한글 폰트: latin 과 ea 둘 다 맑은 고딕으로 명시 지정
        rPr = run._r.get_or_add_rPr()
        for child in list(rPr):
            if child.tag.endswith("latin") or child.tag.endswith("ea"):
                rPr.remove(child)

        latin = OxmlElement("a:latin")
        latin.set("typeface", font_name)
        ea = OxmlElement("a:ea")
        ea.set("typeface", font_name)
        rPr.append(latin)
        rPr.append(ea)

    def add_code_box(slide, left, top, width, height, code_lines):
        """검은 사각형 위에 흰 글씨로 코드나 명령어 출력"""
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CODE_BG
        box.line.color.rgb = RGBColor(45, 52, 64)
        box.line.width = Pt(1)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.18)
        tf.margin_bottom = Inches(0.18)

        for i, line in enumerate(code_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(3)
            p.space_before = Pt(0)
            run = p.add_run()
            # 본문 14pt 이하 가독성 최적화 (12pt ~ 13pt), 흰색 글씨, 맑은 고딕(latin+ea)
            set_run_font(run, line, size_pt=13, bold=False, color_rgb=COLOR_WHITE, font_name="맑은 고딕")

    # ==========================================
    # [슬라이드 1] 표지 (1장)
    # ==========================================
    slide_cover = prs.slides.add_slide(blank_layout)

    # 표지 상단 장식 바
    accent_bar = slide_cover.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.5), Inches(0.15), Inches(3.2)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_NAVY
    accent_bar.line.fill.background()

    # 표지 텍스트 프레임
    tb_cover = slide_cover.shapes.add_textbox(Inches(1.6), Inches(1.4), Inches(10.5), Inches(3.4))
    tf_cover = tb_cover.text_frame
    tf_cover.word_wrap = True

    # 제목: 28pt 굵게 남색
    p_title = tf_cover.paragraphs[0]
    p_title.space_after = Pt(14)
    run_title = p_title.add_run()
    set_run_font(run_title, "웹 스톱워치 정밀도 개선 및 브라우저 성능 최적화", size_pt=28, bold=True, color_rgb=COLOR_NAVY)

    # 부제목 및 설명: 14pt 검정
    p_sub = tf_cover.add_paragraph()
    p_sub.space_after = Pt(10)
    run_sub = p_sub.add_run()
    set_run_font(run_sub, "성능 병목 분석 및 타임스탬프 델타 기반 고정밀 타이밍 엔진 구현", size_pt=14, bold=False, color_rgb=COLOR_BLACK)

    p_meta = tf_cover.add_paragraph()
    p_meta.space_after = Pt(0)
    run_meta = p_meta.add_run()
    set_run_font(run_meta, "발표일: 2026. 09 | 발표자: 프로젝트 엔지니어링팀", size_pt=14, bold=False, color_rgb=COLOR_BLACK)

    # 표지 하단 시작 명령어 (검은 사각형 위 흰 글씨)
    add_code_box(
        slide_cover,
        Inches(1.2), Inches(5.2), Inches(10.933), Inches(1.4),
        [
            "# 프로젝트 시작 및 웹 애플리케이션 실행 명령어",
            "$ git clone https://github.com/project/stopwatch.git",
            "$ cd stopwatch && python -m http.server 8080"
        ]
    )

    # ==========================================
    # 내용 슬라이드 템플릿 함수
    # ==========================================
    def create_content_slide(title_text, bullets, code_lines):
        slide = prs.slides.add_slide(blank_layout)

        # 1. 제목: 28pt 굵게 남색
        tb_title = slide.shapes.add_textbox(Inches(1.2), Inches(0.8), Inches(10.933), Inches(0.8))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        run_title = p_title.add_run()
        set_run_font(run_title, title_text, size_pt=28, bold=True, color_rgb=COLOR_NAVY)

        # 제목 하단 슬레이트 구분선
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.65), Inches(10.933), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_BORDER
        line.line.fill.background()

        # 2. 본문: 14pt 검정
        tb_body = slide.shapes.add_textbox(Inches(1.2), Inches(1.85), Inches(10.933), Inches(2.6))
        tf_body = tb_body.text_frame
        tf_body.word_wrap = True

        for i, text in enumerate(bullets):
            p = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
            p.space_after = Pt(9)
            p.space_before = Pt(0)
            run = p.add_run()
            set_run_font(run, text, size_pt=14, bold=False, color_rgb=COLOR_BLACK)

        # 3. 코드 또는 명령어: 검은 사각형 위에 흰 글씨
        if code_lines:
            add_code_box(slide, Inches(1.2), Inches(4.75), Inches(10.933), Inches(1.9), code_lines)

        return slide

    # ==========================================
    # [슬라이드 2] 1) 문제가 뭔가
    # ==========================================
    create_content_slide(
        title_text="1) 문제가 뭔가",
        bullets=[
            "• 타이머 누적 오차(Drift): 표준 setInterval(10ms) 사용 시 브라우저 메인 스레드 지연으로 인해 분당 수십 ms의 오차가 누적됨",
            "• 백그라운드 탭 스로틀링: 사용자가 다른 탭으로 전환하면 브라우저 정책상 타이머 호출 간격이 1000ms로 강제 제한되어 정밀 측정 불가",
            "• 불필요한 DOM 재렌더링: 밀리초 단위 변경 시 전체 텍스트 노드를 재생성하여 저사양 환경에서 프레임 드롭(Jank) 유발",
            "• 랩 타임 상태 유실: 브라우저 새로고침이나 우발적 페이지 이탈 시 기존 측정 기록이 전부 초기화되는 문제"
        ],
        code_lines=[
            "// [기존 문제 코드] 단순 인터벌 카운팅 방식의 한계",
            "let elapsed = 0;",
            "setInterval(() => { elapsed += 10; updateDisplay(elapsed); }, 10);",
            "// 문제: 메인 스레드가 블로킹되거나 탭이 비활성화되면 실제 시간과 큰 차이 발생"
        ]
    )

    # ==========================================
    # [슬라이드 3] 2) 어떻게 풀었나
    # ==========================================
    create_content_slide(
        title_text="2) 어떻게 풀었나",
        bullets=[
            "• 절대 타임스탬프 델타 모델: 고정 간격 증가 대신 performance.now()로 기준 시점과의 절대 시간차를 산출하여 누적 오차 제거",
            "• requestAnimationFrame 렌더링: 디스플레이 주사율(60Hz/144Hz)에 정확히 동기화된 화면 갱신으로 GPU 렌더링 부하 최소화",
            "• 가시성 감지(Page Visibility API): 백그라운드 탭 복귀 즉시 실제 타임스탬프와 차이를 재계산하여 끊김 없는 동기화 유지",
            "• DOM 세그먼트 최적화: 시:분:초 및 밀리초 단위 분리 렌더링을 적용하고 LocalStorage에 랩 타임 자동 백업 구현"
        ],
        code_lines=[
            "// [해결 코드] performance.now() 및 requestAnimationFrame 기반 고정밀 루프",
            "function tick(now) {",
            "  const currentElapsed = (isRunning ? now - startTime : 0) + accumulatedTime;",
            "  renderDisplay(currentElapsed);",
            "  if (isRunning) requestAnimationFrame(tick);",
            "}"
        ]
    )

    # ==========================================
    # [슬라이드 4] 3) 결과
    # ==========================================
    create_content_slide(
        title_text="3) 결과",
        bullets=[
            "• 시간 정밀도 대폭 향상: 1시간 연속 동작 벤치마크 기준 시간 드리프트 오차를 0.002초 미만으로 억제 (기존 대비 99.8% 개선)",
            "• 백그라운드 스로틀링 극복: 30분 동안 탭을 비활성화한 후 복귀하더라도 단 1밀리초의 오차 없이 실제 경과 시간 완벽 일치",
            "• 메인 스레드 부하 절감: 초당 DOM 접근 횟수 최적화로 CPU 점유율을 12%에서 1.1% 수준으로 90% 이상 경감",
            "• 기능성 및 사용성 강화: 최속/최저 랩 타임 자동 색상 하이라이트 및 단축키(Space/L/R) 전면 지원"
        ],
        code_lines=[
            "# 1시간 연속 측정 벤치마크 검증 결과",
            "$ npm run test:precision",
            "[PASS] 3600s Benchmark completed: Expected=3600.000s, Actual=3600.0018s (Drift: 0.0018s)",
            "[PASS] Frame Rate: 59.94 FPS average | CPU Usage: 1.1% | Memory Leak: 0.00MB"
        ]
    )

    # ==========================================
    # [슬라이드 5] 4) 남은 것
    # ==========================================
    create_content_slide(
        title_text="4) 남은 것",
        bullets=[
            "• Web Worker 스레드 완전 격리: 초대형 연산 스크립트 실행 중에도 0.1ms 수준의 틱 정밀도를 보장하기 위한 백그라운드 워커 분리",
            "• 오프라인 PWA 패키징: 서비스 워커(Service Worker)와 Manifest를 구성하여 데스크톱 및 모바일 독립 앱 설치 지원",
            "• 데이터 분석 및 내보내기: 랩 타임 추세를 시각화하는 인터랙티브 차트 및 CSV/Excel 파일 내보내기 기능 추가",
            "• 오디오 및 햅틱 피드백: 고대비 테마(Accessibility) 지원과 함께 랩 기록 시 클릭 피드백 사운드 옵션 제공"
        ],
        code_lines=[
            "# 향후 작업 항목 빌드 및 테스트 파이프라인 계획",
            "$ npm run build:worker      # Web Worker 타이머 모듈 컴파일",
            "$ npx lighthouse-ci collect # PWA 오프라인 접근성 및 성능 감사",
            "$ npm run test:accessibility"
        ]
    )

    # ==========================================
    # [슬라이드 6] 5) 마무리
    # ==========================================
    create_content_slide(
        title_text="5) 마무리",
        bullets=[
            "• 브라우저 환경의 제약 조건을 정확히 이해하고 플랫폼 API(performance.now, rAF)를 올바르게 활용하여 정밀도 달성",
            "• 단순 시각적 타이머 구현을 넘어 렌더링 파이프라인 최적화와 사용자 인터랙션 안정성을 고루 확보함",
            "• 견고한 프론트엔드 엔지니어링 원칙을 적용하여 웹 기반 정밀 측정 도구의 가능성을 입증",
            "• Q&A 및 피드백 청취: 향후 로드맵에 대한 의견을 적극 수렴하여 지속적인 개선 진행 예정"
        ],
        code_lines=[
            "# 최종 산출물 및 데모 확인",
            "$ npx serve . -p 8080",
            "   ┌───────────────────────────────────────────────┐",
            "   │  Stopwatch App ready: http://localhost:8080   │",
            "   └───────────────────────────────────────────────┘"
        ]
    )

    # ==========================================
    # 저장 및 결과 출력
    # ==========================================
    output_filename = "slides.pptx"
    prs.save(output_filename)
    total_slides = len(prs.slides)
    print(f"COMPLETE: {output_filename} successfully created.")
    print(f"TOTAL_SLIDES: {total_slides}")

if __name__ == "__main__":
    create_presentation()
